"""
PDFやPPTXなどのドキュメントからテキストを抽出するユーティリティ
"""

import os
import re
import io
import PyPDF2
from pptx import Presentation
from docx import Document
from PyQt5.QtCore import QObject, pyqtSignal

class DocumentParser(QObject):
    """ドキュメントからテキストを抽出するクラス"""
    
    progress_updated = pyqtSignal(int, str)
    
    def __init__(self):
        super().__init__()
    
    def extract_text_from_file(self, file_path):
        """
        ファイルからテキストを抽出する
        
        Args:
            file_path (str): ファイルのパス
            
        Returns:
            str: 抽出されたテキスト
        """
        if not os.path.exists(file_path):
            self.progress_updated.emit(100, f"ファイルが見つかりません: {file_path}")
            return ""
            
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.pdf':
                return self.extract_from_pdf(file_path)
            elif file_ext == '.pptx':
                return self.extract_from_pptx(file_path)
            elif file_ext == '.docx':
                return self.extract_from_docx(file_path)
            elif file_ext in ['.txt', '.md', '.csv']:
                return self.extract_from_text(file_path)
            else:
                self.progress_updated.emit(100, f"サポートされていないファイル形式です: {file_ext}")
                return ""
        except Exception as e:
            self.progress_updated.emit(100, f"テキスト抽出エラー: {str(e)}")
            return ""
    
    def extract_from_pdf(self, file_path):
        """
        PDFファイルからテキストを抽出する
        
        Args:
            file_path (str): PDFファイルのパス
            
        Returns:
            str: 抽出されたテキスト
        """
        self.progress_updated.emit(10, "PDFからテキストを抽出中...")
        
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            num_pages = len(reader.pages)
            
            for i, page in enumerate(reader.pages):
                progress = int(10 + (i / num_pages) * 80)
                self.progress_updated.emit(progress, f"PDFページ {i+1}/{num_pages} を処理中...")
                text += page.extract_text() + "\n\n"
                
        self.progress_updated.emit(90, "PDF処理完了")
        return text
    
    def extract_from_pptx(self, file_path):
        """
        PowerPointファイルからテキストを抽出する
        
        Args:
            file_path (str): PowerPointファイルのパス
            
        Returns:
            str: 抽出されたテキスト
        """
        self.progress_updated.emit(10, "PowerPointからテキストを抽出中...")
        
        prs = Presentation(file_path)
        text = ""
        
        total_slides = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            progress = int(10 + (i / total_slides) * 80)
            self.progress_updated.emit(progress, f"スライド {i+1}/{total_slides} を処理中...")
            
            text += f"スライド {i+1}:\n"
            
            if slide.title:
                text += f"タイトル: {slide.title.text}\n"
                
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += f"{shape.text}\n"
            
            text += "\n"
            
        self.progress_updated.emit(90, "PowerPoint処理完了")
        return text
    
    def extract_from_docx(self, file_path):
        """
        Wordファイルからテキストを抽出する
        
        Args:
            file_path (str): Wordファイルのパス
            
        Returns:
            str: 抽出されたテキスト
        """
        self.progress_updated.emit(10, "Word文書からテキストを抽出中...")
        
        doc = Document(file_path)
        text = ""
        
        total_paragraphs = len(doc.paragraphs)
        for i, para in enumerate(doc.paragraphs):
            progress = int(10 + (i / total_paragraphs) * 80)
            self.progress_updated.emit(progress, f"段落 {i+1}/{total_paragraphs} を処理中...")
            
            if para.text:
                text += para.text + "\n"
                
        self.progress_updated.emit(90, "Word文書処理完了")
        return text
    
    def extract_from_text(self, file_path):
        """
        テキストファイルからテキストを抽出する
        
        Args:
            file_path (str): テキストファイルのパス
            
        Returns:
            str: 抽出されたテキスト
        """
        self.progress_updated.emit(10, "テキストファイルを読み込み中...")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                
            self.progress_updated.emit(90, "テキストファイル処理完了")
            return text
        except UnicodeDecodeError:
            # UTF-8でエラーが出た場合は、他のエンコーディングでトライ
            try:
                with open(file_path, 'r', encoding='shift-jis') as f:
                    text = f.read()
                    
                self.progress_updated.emit(90, "テキストファイル処理完了")
                return text
            except:
                self.progress_updated.emit(100, "テキストファイルのエンコーディングを検出できませんでした")
                return "" 